import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

class IngestionAgent:
    def __init__(self):
        pass

    def parse_pdf(self, file):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        return text

    def parse_docx(self, file):
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])

    def parse_pptx(self, file):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def parse_csv(self, file):
        df = pd.read_csv(file)
        return df.to_string(index=False)

    def parse_txt(self, file):
        return file.read().decode("utf-8")

    def chunk_text(self, text, chunk_size=500, overlap=50):
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size]
            chunks.append(chunk.strip())
        return chunks

    def parse_documents(self, files):
        final_chunks = []

        for file in files:
            name = file.name
            if name.endswith(".pdf"):
                text = self.parse_pdf(file)
            elif name.endswith(".docx"):
                text = self.parse_docx(file)
            elif name.endswith(".pptx"):
                text = self.parse_pptx(file)
            elif name.endswith(".csv"):
                text = self.parse_csv(file)
            elif name.endswith(".txt") or name.endswith(".md"):
                text = self.parse_txt(file)
            else:
                continue  # Unsupported file format

            chunks = self.chunk_text(text)
            final_chunks.extend(chunks)

        return final_chunks

    def handle(self, message):
        assert message["type"] == "DOCUMENT_UPLOAD"
        print(f"[MCP][{message['trace_id']}] IngestionAgent received DOCUMENT_UPLOAD")
        files = message["payload"]["files"]
        return self.parse_documents(files)
